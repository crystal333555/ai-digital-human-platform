from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1
s1 = prs.slides.add_slide(prs.slide_layouts[6])
tb = s1.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
p = tb.text_frame.paragraphs[0]
p.text = "AI数字人技术介绍"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = 1

# Slide 2
s2 = prs.slides.add_slide(prs.slide_layouts[6])
tb2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
p2 = tb2.text_frame.paragraphs[0]
p2.text = "核心技术"
p2.font.size = Pt(36)
p2.font.bold = True
for t in ["1. 语音合成技术 (TTS)", "2. 口型同步技术 (Lip Sync)", "3. 面部动画生成"]:
    pp = tb2.text_frame.add_paragraph()
    pp.text = t
    pp.font.size = Pt(24)

# Slide 3
s3 = prs.slides.add_slide(prs.slide_layouts[6])
tb3 = s3.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
p3 = tb3.text_frame.paragraphs[0]
p3.text = "应用场景"
p3.font.size = Pt(36)
p3.font.bold = True
pp2 = tb3.text_frame.add_paragraph()
pp2.text = "虚拟主播、在线教育、企业培训、智能客服"
pp2.font.size = Pt(24)

prs.save("test_ppt.pptx")
