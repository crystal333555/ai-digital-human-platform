from pptx import Presentation
import os

ppt_path = r'D:\AI_Avatar_Data\ppt\c93b043a064040ba8c761726c59a0d0d.pptx'
prs = Presentation(ppt_path)

for i, sl in enumerate(prs.slides):
    texts = []
    for s in sl.shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            texts.append(s.text_frame.text.strip()[:200])
    combined = " | ".join(texts)
    print(f"Slide {i}: {combined[:300]}")
