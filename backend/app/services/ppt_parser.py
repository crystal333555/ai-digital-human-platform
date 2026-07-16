"""PPT解析服务 - 提取PPT页面图片和文字"""

import os
import subprocess
import logging
from typing import List, Dict, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

# 项目根目录
_PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))


def parse_ppt(ppt_path: str, output_dir: str) -> List[Dict]:
    """
    解析PPT文件，提取每页的图片和文字
    
    Args:
        ppt_path: PPT文件绝对路径
        output_dir: 输出目录（存放页面图片）
    
    Returns:
        [{"index": 0, "image_path": "...", "text": "..."}, ...]
    """
    os.makedirs(output_dir, exist_ok=True)
    
    # 优先尝试COM自动化（Windows + PowerPoint已安装时效果最好）
    slides = _try_com_export(ppt_path, output_dir)
    if slides is not None:
        logger.info(f"[PPT] COM export succeeded, {len(slides)} slides")
        return slides
    
    # 降级：python-pptx 渲染
    slides = _parse_with_pptx(ppt_path, output_dir)
    logger.info(f"[PPT] python-pptx parse succeeded, {len(slides)} slides")
    return slides


def _try_com_export(ppt_path: str, output_dir: str) -> Optional[List[Dict]]:
    """尝试用Windows COM自动化导出PPT页面为高清图片（通过PowerShell）"""
    if not os.path.exists(ppt_path):
        return None
    
    abs_ppt = os.path.abspath(ppt_path).replace("'", "''")
    abs_out = os.path.abspath(output_dir).replace("'", "''")
    
    # 使用PowerShell调用PowerPoint COM（比comtypes更稳定）
    ps_script = f"""
if (-not (Test-Path '{abs_out}')) {{ New-Item -ItemType Directory -Path '{abs_out}' -Force | Out-Null }}
$ppt = New-Object -ComObject PowerPoint.Application
$pres = $ppt.Presentations.Open('{abs_ppt}', $true, $false, $false)
$slideCount = $pres.Slides.Count
for ($i = 1; $i -le $slideCount; $i++) {{
    $imgPath = Join-Path '{abs_out}' ("slide_{{0:D3}}.png" -f $i)
    $pres.Slides.Item($i).Export($imgPath, 'PNG', 1920, 1080)
    # Extract text
    $textParts = @()
    $slide = $pres.Slides.Item($i)
    foreach ($shape in $slide.Shapes) {{
        if ($shape.HasTextFrame -eq -1) {{
            try {{
                foreach ($para in $shape.TextFrame.Paragraphs) {{
                    $t = $para.Text.Trim()
                    if ($t) {{ $textParts += $t }}
                }}
            }} catch {{}}
        }}
    }}
    $textFile = Join-Path '{abs_out}' ("slide_{{0:D3}}.txt" -f $i)
    $textParts -join "`n" | Out-File -FilePath $textFile -Encoding utf8
}}
$pres.Close()
$ppt.Quit()
[System.Runtime.Interopservices.Marshal]::ReleaseComObject($ppt) | Out-Null
Write-Host "OK:$slideCount"
"""
    
    try:
        result = subprocess.run(
            ["powershell", "-Command", ps_script],
            capture_output=True, text=True, timeout=120
        )
        
        if "OK:" not in result.stdout:
            logger.warning(f"[PPT] PowerShell COM failed: {result.stderr[:500]}")
            return None
        
        slide_count = int(result.stdout.strip().split(":")[-1])
        slides = []
        
        for i in range(1, slide_count + 1):
            img_path = os.path.join(output_dir, f"slide_{i:03d}.png")
            
            if not os.path.exists(img_path):
                continue
            
            # 始终用python-pptx提取文字（COM提取经常因编码问题失败）
            text = _extract_text_with_pptx(ppt_path, i - 1)
            
            slides.append({
                "index": i - 1,
                "image_path": img_path,
                "text": text,
            })
        
        return slides if slides else None
        
    except subprocess.TimeoutExpired:
        logger.warning("[PPT] PowerShell COM timeout")
        return None
    except Exception as e:
        logger.warning(f"[PPT] PowerShell COM failed: {e}")
        return None


def _extract_text_with_pptx(ppt_path: str, slide_index: int) -> str:
    """用python-pptx提取指定页的文字"""
    try:
        prs = Presentation(ppt_path)
        if slide_index < len(prs.slides):
            text_parts = []
            for shape in prs.slides[slide_index].shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            text_parts.append(t)
            return "\n".join(text_parts)
    except Exception as e:
        logger.warning(f"[PPT] python-pptx text extraction failed: {e}")
    return ""


def _parse_with_pptx(ppt_path: str, output_dir: str) -> List[Dict]:
    """用 python-pptx 解析PPT，提取文字并渲染页面图片"""
    prs = Presentation(ppt_path)
    slides = []
    
    # 获取PPT尺寸（EMU转像素，1 inch = 914400 EMU, 假设96 DPI）
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    width_px = int(slide_width / 914400 * 150)  # 放大150倍获得高清图
    height_px = int(slide_height / 914400 * 150)
    # 限制最大尺寸
    max_dim = 1920
    if width_px > max_dim or height_px > max_dim:
        scale = max_dim / max(width_px, height_px)
        width_px = int(width_px * scale)
        height_px = int(height_px * scale)
    
    for idx, slide in enumerate(prs.slides):
        # 提取文字
        text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        text_parts.append(t)
        
        extracted_text = "\n".join(text_parts)
        
        # 渲染页面图片（简化版：白色背景+文字）
        img_path = os.path.join(output_dir, f"slide_{idx + 1:03d}.png")
        _render_slide_image(slide, width_px, height_px, img_path)
        
        slides.append({
            "index": idx,
            "image_path": img_path,
            "text": extracted_text,
        })
    
    return slides


def _render_slide_image(slide, width: int, height: int, output_path: str):
    """将PPT页面渲染为图片（简化版：背景+文字+形状色块）"""
    img = Image.new("RGB", (width, height), (255, 255, 255))
    draw = ImageDraw.Draw(img)
    
    # 获取slide尺寸用于坐标转换
    try:
        slide_part = slide.part
        # 通过 slide_part 获取 presentation 的尺寸
        prs_part = slide_part.package.presentation_part
        slide_width_emu = prs_part.presentation.slide_width
        slide_height_emu = prs_part.presentation.slide_height
    except Exception:
        # 默认16:9
        slide_width_emu = 12192000  # 13.333 inches
        slide_height_emu = 6858000  # 7.5 inches
    
    def emu_to_px(emu_val, total_emu, total_px):
        return int(emu_val / total_emu * total_px)
    
    # 绘制形状和文字
    try:
        for shape in slide.shapes:
            # 坐标转换
            left = emu_to_px(shape.left, slide_width_emu, width)
            top = emu_to_px(shape.top, slide_height_emu, height)
            w = emu_to_px(shape.width, slide_width_emu, width)
            h = emu_to_px(shape.height, slide_height_emu, height)
            
            # 绘制背景色块
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                try:
                    fill = shape.fill
                    if fill.type is not None and hasattr(fill, 'fore_color') and fill.fore_color is not None:
                        rgb = fill.fore_color.rgb
                        if rgb:
                            color = (rgb[0], rgb[1], rgb[2])
                            draw.rectangle([left, top, left + w, top + h], fill=color)
                except:
                    pass
            
            # 绘制文字
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if not para.text.strip():
                        continue
                    # 计算字体大小
                    font_size = 16
                    try:
                        if para.runs and para.runs[0].font.size:
                            font_size = max(8, int(para.runs[0].font.size / 12700 * width / 960))
                    except:
                        pass
                    
                    # 文字颜色
                    text_color = (0, 0, 0)
                    try:
                        if para.runs and para.runs[0].font.color and para.runs[0].font.color.rgb:
                            rgb = para.runs[0].font.color.rgb
                            text_color = (rgb[0], rgb[1], rgb[2])
                    except:
                        pass
                    
                    # 使用默认字体
                    try:
                        font = ImageFont.truetype("arial.ttf", font_size)
                    except:
                        font = ImageFont.load_default()
                    
                    draw.text((left + 4, top + 2), para.text, fill=text_color, font=font)
                    top += font_size + 4
    except Exception as e:
        logger.warning(f"[PPT] Render slide error: {e}")
    
    img.save(output_path, "PNG")
